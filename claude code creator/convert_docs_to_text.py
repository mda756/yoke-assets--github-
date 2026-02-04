"""
DOCX/PPTX to TXT Batch Converter
Recursively converts all .docx and .pptx files to .txt in a folder
"""

import os
from pathlib import Path

def convert_docx_to_txt(docx_path):
    """Convert .docx to .txt using python-docx"""
    try:
        from docx import Document
        doc = Document(docx_path)
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return '\n'.join(text)
    except Exception as e:
        return f"ERROR converting {docx_path}: {str(e)}"

def convert_pptx_to_txt(pptx_path):
    """Convert .pptx to .txt using python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"\n--- SLIDE {slide_num} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        return f"ERROR converting {pptx_path}: {str(e)}"

def batch_convert(root_folder):
    """Convert all .docx and .pptx files in folder and subfolders"""
    root_path = Path(root_folder)

    # Find all .docx files
    docx_files = list(root_path.rglob("*.docx"))
    pptx_files = list(root_path.rglob("*.pptx"))

    print(f"Found {len(docx_files)} .docx files and {len(pptx_files)} .pptx files")

    converted = 0
    errors = 0

    # Convert DOCX files
    for docx_file in docx_files:
        # Skip temp files
        if docx_file.name.startswith('~$'):
            continue

        txt_file = docx_file.with_suffix('.txt')

        # Skip if already converted
        if txt_file.exists():
            print(f"SKIP (exists): {docx_file.name}")
            continue

        print(f"Converting: {docx_file.name}")
        text_content = convert_docx_to_txt(docx_file)

        if text_content.startswith("ERROR"):
            print(f"  ERROR: {text_content}")
            errors += 1
        else:
            txt_file.write_text(text_content, encoding='utf-8')
            print(f"  OK: Created {txt_file.name}")
            converted += 1

    # Convert PPTX files
    for pptx_file in pptx_files:
        # Skip temp files
        if pptx_file.name.startswith('~$'):
            continue

        txt_file = pptx_file.with_suffix('.txt')

        # Skip if already converted
        if txt_file.exists():
            print(f"SKIP (exists): {pptx_file.name}")
            continue

        print(f"Converting: {pptx_file.name}")
        text_content = convert_pptx_to_txt(pptx_file)

        if text_content.startswith("ERROR"):
            print(f"  ERROR: {text_content}")
            errors += 1
        else:
            txt_file.write_text(text_content, encoding='utf-8')
            print(f"  OK: Created {txt_file.name}")
            converted += 1

    print(f"\nDONE: {converted} files converted, {errors} errors")

if __name__ == "__main__":
    # Default folder - change this or pass as argument
    folder = r"C:\Users\matth\Dropbox\Yoke Digital\yoke-assets--github-\chatgpt instructions and knowledge (For ingest by cluade code)\MASTER KOWELDGE PACK"

    print(f"Converting all .docx and .pptx files in:\n{folder}\n")

    # Check if libraries are installed
    try:
        import docx
        import pptx
    except ImportError:
        print("ERROR: Required libraries not installed!")
        print("\nRun these commands first:")
        print("  pip install python-docx")
        print("  pip install python-pptx")
        exit(1)

    batch_convert(folder)
